import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Dark Cyberpunk / Sleek Enterprise FinOps)
BG_DARK = RGBColor(11, 15, 25)        # #0B0F19
CARD_BG = RGBColor(22, 30, 48)        # #161E30
CARD_BORDER = RGBColor(45, 60, 90)    # #2D3C5A
TEXT_LIGHT = RGBColor(240, 244, 248)  # White/Off-white
TEXT_MUTED = RGBColor(148, 163, 184)  # Slate gray
ACCENT_BLUE = RGBColor(56, 189, 248)  # Sky Blue
ACCENT_PURPLE = RGBColor(168, 85, 247)# Vibrant Purple
ACCENT_CYAN = RGBColor(45, 212, 191)   # Teal Cyan
ACCENT_GREEN = RGBColor(52, 211, 153)  # Mint Green
ACCENT_AMBER = RGBColor(251, 191, 36)  # Amber Yellow
ACCENT_ROSE = RGBColor(244, 63, 94)    # Rose Red

def add_blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return slide

def add_header(slide, category, title, subtitle=None):
    tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.35))
    tf_cat = tb_cat.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_CYAN
    p_cat.font.name = 'Calibri'
    
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.7))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT
    p_title.font.name = 'Calibri'

    if subtitle:
        p_sub = tf_title.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.font.name = 'Calibri'

def add_card(slide, left, top, width, height, border_color=CARD_BORDER, fill_color=CARD_BG):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    return card

# ====================================================
# SLIDE 1: Title Slide & Project Overview
# ====================================================
s1 = add_blank_slide()
add_card(s1, Inches(1.2), Inches(1.0), Inches(10.933), Inches(5.5), border_color=ACCENT_PURPLE)

tb1 = s1.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "CLOUD ATLAS AI"
p.font.size = Pt(46)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN
p.alignment = PP_ALIGN.CENTER

p2 = tf1.add_paragraph()
p2.text = "Next-Generation Multi-Cloud FinOps & Predictive Governance Platform"
p2.font.size = Pt(22)
p2.font.color.rgb = TEXT_LIGHT
p2.alignment = PP_ALIGN.CENTER

p3 = tf1.add_paragraph()
p3.text = "\nPredictIQ 90-Day ML Forecasting | One-Class SVM Anomaly Shield | Workload Simulator | Migration Arbitrage"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED
p3.alignment = PP_ALIGN.CENTER

p4 = tf1.add_paragraph()
p4.text = "\nPresented by: Neel Panchal & Henil  |  Tech Stack: React 19, Node.js, Express, MongoDB Atlas & Python ML"
p4.font.size = Pt(14)
p4.font.bold = True
p4.font.color.rgb = ACCENT_BLUE
p4.alignment = PP_ALIGN.CENTER


# ====================================================
# SLIDE 2: Industry Problem & Executive Vision
# ====================================================
s2 = add_blank_slide()
add_header(s2, "Problem & Vision", "Enterprise FinOps Crisis & The CloudAtlas Solution", "Solving multi-cloud visibility gaps, unexpected budget spikes, and static reporting.")

cards_s2 = [
    ("💸 Multi-Cloud Billing Opacity", "Raw billing files across AWS, Azure, and GCP contain millions of unaggregated rows. FinOps teams lack clear unit cost visibility.", ACCENT_AMBER),
    ("🚨 Runaway Spikes & Leaks", "Orphan EBS volumes, unattached Elastic IPs, and runaway serverless functions silently cause massive budget overruns.", ACCENT_ROSE),
    ("🤖 Autonomous AI Solution", "CloudAtlas AI replaces static historical graphs with ML engines that predict future spend 90 days out and auto-detect anomalies.", ACCENT_CYAN)
]

for i, (title, text, color) in enumerate(cards_s2):
    left = Inches(0.8 + i * 3.95)
    add_card(s2, left, Inches(1.8), Inches(3.7), Inches(4.8), border_color=color)
    tb = s2.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    
    p_body = tf.add_paragraph()
    p_body.text = f"\n{text}"
    p_body.font.size = Pt(13.5)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 3: System Architecture & Full-Stack Tech Stack
# ====================================================
s3 = add_blank_slide()
add_header(s3, "Architecture & Stack", "3-Tier Microservices Infrastructure & Real-Time Data Pipeline", "High-performance microservices, zero-trust security, and real-time state synchronization.")

stacks = [
    ("Frontend Layer (`/client`)", "• React 19 + Vite 6\n• Framer Motion & GSAP ScrollTrigger\n• Recharts Analytics Visualization\n• Tailwind CSS Dark Cyberpunk Theme", ACCENT_CYAN, Inches(0.8)),
    ("Backend API Gateway (`/server`)", "• Node.js & Express REST Gateway\n• Zero-Trust JWT Security Guards\n• Multer CSV Ingestion Pipeline\n• Express Rate-Limiter & SOC2 Logger", ACCENT_BLUE, Inches(4.8)),
    ("ML Engine & MongoDB Atlas", "• Python ML (PredictIQ & SVM Shield)\n• MongoDB Mongoose ORM Schemas\n• BillingData & UploadedFile Models\n• Real-Time Signal Bus (DataContext)", ACCENT_PURPLE, Inches(8.8))
]

for title, desc, color, left in stacks:
    add_card(s3, left, Inches(1.8), Inches(3.7), Inches(4.8), border_color=color)
    tb = s3.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p_body = tf.add_paragraph()
    p_body.text = f"\n{desc}"
    p_body.font.size = Pt(12.5)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 4: Core Features Blueprint (All Modules In One)
# ====================================================
s4 = add_blank_slide()
add_header(s4, "Platform Blueprint", "Comprehensive 14+ Module Ecosystem Overview", "Unified control center for executive dashboards, AI engines, and security governance.")

modules = [
    ("📊 Executive Dashboard", "Aggregated total spend, provider distribution (AWS vs Azure vs GCP), and top expensive services."),
    ("🔮 PredictIQ 90-Day Forecast", "XGBoost time-series expenditure forecasting integrating linear trends and Fourier seasonal components."),
    ("🛡️ One-Class SVM Anomaly Shield", "Unsupervised anomaly detection flagging billing spikes (z >= 1.6) and orphan resource leakage."),
    ("⚡ Interactive What-If Simulator", "Real-time scale testing (CPU %, RAM, Storage, Egress) to evaluate Reserved Instance & Cloud migration savings."),
    ("🌐 Migration Intelligence", "Cross-cloud workload cost arbitrage comparing AWS EC2, Azure VMs, and GCP Compute Engine."),
    ("🔐 Gatekeeper & SOC2 Audit", "Double-lock Super Admin portal (`/admin/login` passcode: `HenilNeelProject`) and audit logging.")
]

for i, (title, desc) in enumerate(modules):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 5.9)
    top = Inches(1.8 + row * 1.7)
    
    add_card(s4, left, top, Inches(5.6), Inches(1.5), border_color=CARD_BORDER)
    tb = s4.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), Inches(5.3), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    p_body = tf.add_paragraph()
    p_body.text = desc
    p_body.font.size = Pt(11.5)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 5: Machine Learning Core (PredictIQ & Anomaly Shield)
# ====================================================
s5 = add_blank_slide()
add_header(s5, "AI & Machine Learning", "PredictIQ Forecast & One-Class SVM Anomaly Shield", "Mathematical algorithms driving high-precision cost projections and threat detection.")

add_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_PURPLE)
tb = s5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🔮 1. PredictIQ XGBoost Regressor"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_PURPLE
p_body = tf.add_paragraph()
p_body.text = "\n• Objective: 90-Day time-series daily cost forecasting.\n• Fourier Seasonality Formula:\n  y_hat = Trend(t) + sin(t * 0.4) * 0.09 * mean_cost\n• Multi-Variable Scaling: Adjusts predictions dynamically based on CPU %, RAM, Storage, and Egress sliders."
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = TEXT_LIGHT

add_card(s5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), border_color=ACCENT_CYAN)
tb2 = s5.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
tf2 = tb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "🛡️ 2. One-Class SVM Anomaly Shield"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = ACCENT_CYAN
p_body2 = tf2.add_paragraph()
p_body2.text = "\n• Unsupervised Learning: Fits a tight decision boundary hypersphere around normal baseline operational spend.\n• Z-Score Threshold: Flags records with z >= 1.6 standard deviations above historical mean.\n• Remediation: 1-Click action commands to delete orphan EBS storage or release idle IPs."
p_body2.font.size = Pt(12.5)
p_body2.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 6: FinOps Optimization, Security & Migration
# ====================================================
s6 = add_blank_slide()
add_header(s6, "Governance & Migration", "FinOps Savings, Security Gatekeeper & Migration Arbitrage", "Enterprise protection, cost reduction strategies, and multi-cloud migration engine.")

gov_cards = [
    ("📦 FinOps Savings Engine", "Identifies orphan EBS volumes, idle load balancers, and recommends Reserved Instance (RI) conversion saving up to 60%.", ACCENT_GREEN),
    ("🔐 Zero-Trust Security & Gatekeeper", "JWT authentication, Express rate-limiting, SOC2 logging, and Super Admin portal (`/admin/login` passcode: `HenilNeelProject`).", ACCENT_ROSE),
    ("🌐 Migration Arbitrage Engine", "Compares workload unit pricing across AWS EC2, Azure VMs, and GCP Compute Engine with carbon emission ratings.", ACCENT_BLUE)
]

for i, (title, text, color) in enumerate(gov_cards):
    left = Inches(0.8 + i * 3.95)
    add_card(s6, left, Inches(1.8), Inches(3.7), Inches(4.8), border_color=color)
    tb = s6.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = color
    
    p_body = tf.add_paragraph()
    p_body.text = f"\n{text}"
    p_body.font.size = Pt(13)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 7: Business ROI, Deployment & Conclusion
# ====================================================
s7 = add_blank_slide()
add_header(s7, "Business ROI & Conclusion", "Quantifiable Impact, Operational Deployment & Roadmap", "Measurable achievements, simple startup guide, and conclusion.")

add_card(s7, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.8), border_color=ACCENT_GREEN)
tb = s7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.3), Inches(4.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "30 - 40%"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
p_sub = tf.add_paragraph()
p_sub.text = "Annual Cost Savings\n\n• < 2 Mins Anomaly Detection Speed.\n• 95%+ PredictIQ Forecast Accuracy."
p_sub.font.size = Pt(13)
p_sub.font.color.rgb = TEXT_LIGHT

add_card(s7, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.8), border_color=ACCENT_CYAN)
tb2 = s7.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.3), Inches(4.4))
tf2 = tb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "⚡ Quick Startup"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = ACCENT_CYAN
p_sub2 = tf2.add_paragraph()
p_sub2.text = "\nBackend (`/server`):\n`npm run dev` (Port 5000)\n\nFrontend (`/client`):\n`npm run dev` (Port 5173)\n\nDB: MongoDB Atlas"
p_sub2.font.size = Pt(12.5)
p_sub2.font.color.rgb = TEXT_LIGHT

add_card(s7, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), border_color=ACCENT_PURPLE)
tb3 = s7.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.3), Inches(4.4))
tf3 = tb3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "🚀 Future Roadmap"
p3.font.size = Pt(18)
p3.font.bold = True
p3.font.color.rgb = ACCENT_PURPLE
p_sub3 = tf3.add_paragraph()
p_sub3.text = "\n• CloudWatch Auto-Remediation.\n• Kubernetes Pod FinOps (Prometheus).\n• LLM Natural Language FinOps Copilot.\n\nThank You! Questions & Q&A."
p_sub3.font.size = Pt(12.5)
p_sub3.font.color.rgb = TEXT_LIGHT

prs.save("CloudAtlas_AI_7_Slide_Presentation.pptx")
print("Successfully generated exact 7-slide master presentation: CloudAtlas_AI_7_Slide_Presentation.pptx")
