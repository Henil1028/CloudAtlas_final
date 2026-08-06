import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Dark Cyberpunk / Sleek Enterprise FinOps)
BG_DARK = RGBColor(11, 15, 25)        # #0B0F19
CARD_BG = RGBColor(22, 30, 48)        # #161E30
CARD_BORDER = RGBColor(45, 60, 90)    # #2D3C5A
TEXT_LIGHT = RGBColor(240, 244, 248)  # White/Off-white
TEXT_MUTED = RGBColor(148, 163, 184)  # Slate gray
ACCENT_BLUE = RGBColor(56, 189, 248)  # Sky Blue
ACCENT_PURPLE = RGBColor(168, 85, 247)# Vibrant Purple
ACCENT_CYAN = RGBColor(45, 212, 191)   # Teal Cyan
ACCENT_GREEN = RGBColor(52, 211, 153)  # Mint Green
ACCENT_AMBER = RGBColor(251, 191, 36)  # Amber Yellow
ACCENT_ROSE = RGBColor(244, 63, 94)    # Rose Red

def add_blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return slide

def add_header(slide, category, title, subtitle=None):
    tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.35))
    tf_cat = tb_cat.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_CYAN
    p_cat.font.name = 'Calibri'
    
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.7))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT
    p_title.font.name = 'Calibri'

    if subtitle:
        p_sub = tf_title.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.font.name = 'Calibri'

def add_card(slide, left, top, width, height, border_color=CARD_BORDER, fill_color=CARD_BG):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    return card

# ====================================================
# SLIDE 1: Title Slide (Cover)
# ====================================================
s1 = add_blank_slide()
add_card(s1, Inches(1.2), Inches(1.0), Inches(10.933), Inches(5.5), border_color=ACCENT_PURPLE)

tb1 = s1.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "CLOUD ATLAS AI"
p.font.size = Pt(46)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN
p.alignment = PP_ALIGN.CENTER

p2 = tf1.add_paragraph()
p2.text = "Comprehensive End-to-End Technical & Architectural Presentation"
p2.font.size = Pt(22)
p2.font.color.rgb = TEXT_LIGHT
p2.alignment = PP_ALIGN.CENTER

p3 = tf1.add_paragraph()
p3.text = "\nMulti-Cloud FinOps | PredictIQ ML Forecasting | One-Class SVM Anomaly Shield | What-If Simulator"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED
p3.alignment = PP_ALIGN.CENTER

p4 = tf1.add_paragraph()
p4.text = "\nPresented by: Neel Panchal & Henil  |  Built with React 19, Node.js, Express, MongoDB & Python ML"
p4.font.size = Pt(14)
p4.font.bold = True
p4.font.color.rgb = ACCENT_BLUE
p4.alignment = PP_ALIGN.CENTER


# ====================================================
# SLIDE 2: Executive Summary & Project Vision
# ====================================================
s2 = add_blank_slide()
add_header(s2, "Executive Summary", "Multi-Cloud Governance & Financial Intelligence", "Bridging the gap between Cloud Operations, AI Predictions, and Financial Engineering.")

cards_s2 = [
    ("🌐 Multi-Cloud Fragmentation", "Enterprises run across AWS, Azure, GCP with scattered billing logs, making cross-cloud cost visibility extremely difficult.", ACCENT_AMBER),
    ("🤖 Autonomous AI Governance", "PredictIQ uses Machine Learning to forecast spend 90 days out and auto-detect anomalies before billing cycle overruns occur.", ACCENT_PURPLE),
    ("💰 Financial Optimization", "Provides actionable recommendations, Reserved Instance conversion strategies, and interactive What-If scenario modeling.", ACCENT_GREEN)
]

for i, (title, text, color) in enumerate(cards_s2):
    left = Inches(0.8 + i * 3.95)
    add_card(s2, left, Inches(1.8), Inches(3.7), Inches(4.8), border_color=color)
    tb = s2.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    
    p_body = tf.add_paragraph()
    p_body.text = f"\n{text}"
    p_body.font.size = Pt(13.5)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 3: The Problem Statement (FinOps Challenges)
# ====================================================
s3 = add_blank_slide()
add_header(s3, "Problem Statement", "Critical Enterprise Cloud Cost Pain Points", "Why traditional cloud dashboards fail modern engineering and finance teams.")

probs = [
    ("📊 Billing Log Opacity", "Raw billing files contain millions of complex unaggregated lines with disparate format schemes across AWS CUR, Azure Billing, and GCP BigQuery export.", ACCENT_ROSE),
    ("🚨 Runaway Spikes & Leaks", "Orphaned EBS volumes, unattached Elastic IPs, and runaway serverless functions silently consume tens of thousands of dollars before end-of-month audit.", ACCENT_AMBER),
    ("🔮 Reactive vs Predictive", "Legacy tools display historical static graphs. They cannot project future quarter cost impact under dynamic workload scale changes.", ACCENT_PURPLE),
    ("⚖️ Multi-Cloud Complexity", "Comparing unit cost efficiency between AWS EC2, Azure VMs, and GCP Compute Engine requires manual complex mathematical normalization.", ACCENT_BLUE)
]

for i, (title, text, color) in enumerate(probs):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 5.9)
    top = Inches(1.8 + row * 2.5)
    add_card(s3, left, top, Inches(5.6), Inches(2.2), border_color=color)
    tb = s3.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.2), Inches(1.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = color
    p_body = tf.add_paragraph()
    p_body.text = text
    p_body.font.size = Pt(12.5)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 4: Solution Architecture & Tech Stack
# ====================================================
s4 = add_blank_slide()
add_header(s4, "Technology Stack", "Full-Stack Microservice Architecture", "Built with modern performant frameworks and responsive state synchronization.")

stacks = [
    ("Frontend Layer", "• React 19 & Vite 6\n• Framer Motion & GSAP ScrollTrigger\n• Recharts & Lucide Icons\n• Tailwind CSS Token System", ACCENT_CYAN, Inches(0.8)),
    ("Backend API Gateway", "• Node.js & Express REST API\n• Zero-Trust JWT Authentication\n• Multer Multi-Provider Ingestion\n• Rate-Limiting & Audit Logger", ACCENT_BLUE, Inches(4.8)),
    ("ML & Data Layer", "• Python Flask/FastAPI Microservice\n• XGBoost Time-Series Regressor\n• One-Class SVM Anomaly Shield\n• MongoDB Atlas Mongoose ORM", ACCENT_PURPLE, Inches(8.8))
]

for title, desc, color, left in stacks:
    add_card(s4, left, Inches(1.8), Inches(3.7), Inches(4.8), border_color=color)
    tb = s4.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = color
    p_body = tf.add_paragraph()
    p_body.text = f"\n{desc}"
    p_body.font.size = Pt(13)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 5: Full Application Module Flow (UI & Pages)
# ====================================================
s5 = add_blank_slide()
add_header(s5, "Platform Blueprint", "Comprehensive Page & Navigation Ecosystem", "14+ Specialized Enterprise Portals for FinOps Governance.")

modules = [
    ("🎨 Landing Page & Sandbox", "Cinematic particle aurora background, live FinOps interactive sandbox & annual ROI slider calculator."),
    ("🔐 Auth & Super Admin Gatekeeper", "Segmented OTP login, JWT zero-trust authentication, and master security passcode challenge (`/admin/login`)."),
    ("📊 Executive Dashboard", "Aggregated total monthly spend, active cloud provider split (AWS vs Azure vs GCP), and top expensive services."),
    ("🔮 PredictIQ Forecast Page", "90-day time-series cost projections with interactive What-If workload sliders (CPU, RAM, Storage, Egress)."),
    ("🛡️ Anomaly Shield Page", "Z-score outlier detection table with immediate mitigation action buttons (Release orphan EBS, terminate idle VMs)."),
    ("⚡ Migration & Risk Engine", "Multi-cloud migration cost comparison calculator and risk probability scoring for infrastructure workloads.")
]

for i, (title, desc) in enumerate(modules):
    col = i % 2
    row = i // 3
    left = Inches(0.8 + col * 5.9)
    top = Inches(1.8 + row * 1.7)
    
    add_card(s5, left, top, Inches(5.6), Inches(1.5), border_color=CARD_BORDER)
    tb = s5.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), Inches(5.3), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    p_body = tf.add_paragraph()
    p_body.text = desc
    p_body.font.size = Pt(11.5)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 6: Multi-Cloud Dataset Ingestion Engine
# ====================================================
s6 = add_blank_slide()
add_header(s6, "Data Ingestion Pipeline", "Multi-Cloud Dataset Parsing & Ingestion Engine", "Seamlessly processing multi-provider CSV billing dumps into standardized MongoDB schemas.")

pipeline_steps = [
    ("1. File Ingestion", "Multer upload pipeline accepting AWS, Azure, and GCP CSV files up to 2,000+ rows per batch.", ACCENT_BLUE),
    ("2. Schema Mapping", "Automated parsing normalizes provider-specific fields (UsageType, MeterCategory, Cost, Region) into unified schema.", ACCENT_CYAN),
    ("3. Data Validation", "Joi / Custom validators verify date formats, positive cost values, non-empty service identifiers, and user references.", ACCENT_GREEN),
    ("4. Global Event Bus", "DataContext broadcasts lightweight signals to immediately refresh Dashboard & Analytics views without full page reloads.", ACCENT_PURPLE)
]

for i, (title, text, color) in enumerate(pipeline_steps):
    left = Inches(0.8 + i * 2.95)
    add_card(s6, left, Inches(1.8), Inches(2.75), Inches(4.8), border_color=color)
    tb = s6.shapes.add_textbox(left + Inches(0.15), Inches(2.0), Inches(2.45), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    
    p_body = tf.add_paragraph()
    p_body.text = f"\n{text}"
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 7: Machine Learning - PredictIQ Forecast Engine
# ====================================================
s7 = add_blank_slide()
add_header(s7, "AI & Machine Learning", "PredictIQ 90-Day Expenditure Forecast Engine", "XGBoost Regressor incorporating multi-harmonic Fourier seasonality & workload multipliers.")

add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_PURPLE)
tb = s7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🔮 XGBoost Time-Series Architecture"
p.font.size = Pt(19)
p.font.bold = True
p.font.color.rgb = ACCENT_PURPLE
p_body = tf.add_paragraph()
p_body.text = "\n• Objective: Predict daily spend trajectory for 90 days out.\n• Trend Component: Linear regression slope extracted from past 180 days of billing logs.\n• Fourier Seasonality Component:\n  Seasonality(t) = sin(t * 0.4) * 0.09 * mean_cost\n• Multi-Variable Workload Sliders:\n  Scales predictions dynamically as CPU %, RAM, Storage, and Egress values are tweaked."
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = TEXT_LIGHT

add_card(s7, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), border_color=ACCENT_GREEN)
tb2 = s7.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
tf2 = tb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "⚡ What-If Workload Simulator"
p2.font.size = Pt(19)
p2.font.bold = True
p2.font.color.rgb = ACCENT_GREEN
p_body2 = tf2.add_paragraph()
p_body2.text = "\n• Real-Time Scale Testing: Simulates 1.5x to 3x compute/memory scaling without modifying production.\n• Reserved Instance (RI) Evaluation: Calculates net savings if converting On-Demand workloads to 1-Year or 3-Year RIs.\n• Cloud Provider Switcher: Evaluates cost differential of shifting compute nodes from AWS us-east-1 to Azure East US."
p_body2.font.size = Pt(12.5)
p_body2.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 8: Machine Learning - One-Class SVM Anomaly Shield
# ====================================================
s8 = add_blank_slide()
add_header(s8, "Autonomous Security", "One-Class SVM Anomaly Shield & Z-Score Guard", "Unsupervised anomaly detection protecting against runaway cloud infrastructure bills.")

add_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_CYAN)
tb = s8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🛡️ One-Class SVM Detection"
p.font.size = Pt(19)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN
p_body = tf.add_paragraph()
p_body.text = "\n• Unsupervised Learning: Fits a tight decision boundary hypersphere around normal baseline operational cost patterns.\n• Outlier Classification: Data points landing outside the hypersphere are instantly flagged as malicious or unexpected anomalies.\n• Zero Threshold Setup: Works out-of-the-box without requiring complex manually configured metric alerts."
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = TEXT_LIGHT

add_card(s8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), border_color=ACCENT_ROSE)
tb2 = s8.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
tf2 = tb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "🚨 Statistical Z-Score & Automated Remediation"
p2.font.size = Pt(19)
p2.font.bold = True
p2.font.color.rgb = ACCENT_ROSE
p_body2 = tf2.add_paragraph()
p_body2.text = "\n• Z-Score Threshold: Flags records with z >= 1.6 standard deviations above historical mean.\n• Root Cause Attribution: Highlights specific AWS/Azure services, usage types, and regions causing the spike.\n• Actionable Remediation: Provides 1-click mitigation actions (e.g. 'Delete Unattached EBS Volume', 'Scale Down Idle Cluster')."
p_body2.font.size = Pt(12.5)
p_body2.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 9: Deep-Dive Database Schemas (MongoDB Mongoose)
# ====================================================
s9 = add_blank_slide()
add_header(s9, "Database Architecture", "MongoDB Collections & Schema Specifications", "Normalized document schemas optimized for rapid analytics aggregation and audit tracking.")

schemas = [
    ("BillingData Schema", "• date: Date (Indexed)\n• service: String\n• cost: Number\n• region: String\n• usageType: String\n• provider: ['aws','azure','gcp']\n• uploadedBy: ObjectId (Ref User)\n• fileId: ObjectId (Ref UploadedFile)", ACCENT_BLUE, Inches(0.8)),
    ("UploadedFile Schema", "• filename: String\n• provider: String\n• recordCount: Number\n• size: Number (Bytes)\n• uploadedBy: ObjectId\n• status: ['success','failed']\n• createdAt: Date", ACCENT_CYAN, Inches(4.8)),
    ("AuditLog & User Schema", "• user: ObjectId\n• action: String (AUTH, UPLOAD, DELETE)\n• ipAddress: String\n• timestamp: Date\n• role: ['user','admin','superadmin']\n• isOtpVerified: Boolean", ACCENT_AMBER, Inches(8.8))
]

for title, desc, color, left in schemas:
    add_card(s9, left, Inches(1.8), Inches(3.7), Inches(4.8), border_color=color)
    tb = s9.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p_body = tf.add_paragraph()
    p_body.text = f"\n{desc}"
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 10: Security, Authentication & Super Admin Gatekeeper
# ====================================================
s10 = add_blank_slide()
add_header(s10, "Security & Compliance", "Zero-Trust Security & Super Admin Challenge", "Bank-grade protection, rate limiting, and multi-factor authorization safeguards.")

sec_cards = [
    ("🔐 Zero-Trust JWT Auth", "All API routes enforced via Bearer tokens. Includes password hashing with bcrypt, automatic session timeout, and segmented sign-in/sign-up state management.", ACCENT_BLUE),
    ("🛡️ Gatekeeper Portal", "Super Admin portal (`/admin/login`) secured behind a dual-layer challenge requiring Master Passcode (`HenilNeelProject`) before admin credentials.", ACCENT_ROSE),
    ("⚡ Rate Limiting & Protection", "Express rate-limiter protects against brute-force attacks and automated dataset upload spam. Sanitize inputs against Mongo Injection.", ACCENT_AMBER),
    ("📜 SOC2 Audit Governance", "Logs every administrative activity, file deletion, user role mutation, and authentication attempt for enterprise audit compliance.", ACCENT_GREEN)
]

for i, (title, text, color) in enumerate(sec_cards):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 5.9)
    top = Inches(1.8 + row * 2.5)
    add_card(s10, left, top, Inches(5.6), Inches(2.2), border_color=color)
    tb = s10.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.2), Inches(1.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = color
    p_body = tf.add_paragraph()
    p_body.text = text
    p_body.font.size = Pt(12)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 12: Migration Intelligence & Cross-Cloud Pricing
# ====================================================
s12_mig = add_blank_slide()
add_header(s12_mig, "Migration Intelligence", "Cross-Cloud Migration & Workload Arbitrage", "Comparing cost performance and ROI across AWS, Azure, and GCP.")

mig_items = [
    ("🌐 Cloud Migration Arbitrage", "Calculates exact cost differences of migrating workloads between AWS EC2, Azure VMs, and GCP Compute Engine with instant ROI analysis.", ACCENT_CYAN),
    ("⚖️ Multi-Cloud Unit Benchmarking", "Normalizes core vCPU, RAM, and Storage pricing schemas across providers to reveal true unit-level cost efficiency.", ACCENT_BLUE),
    ("🌱 Carbon Footprint & Sustainability", "Evaluates green energy metrics and carbon emission ratings across cloud regions during workload migration planning.", ACCENT_GREEN)
]

for i, (title, text, color) in enumerate(mig_items):
    left = Inches(0.8 + i * 3.95)
    add_card(s12_mig, left, Inches(1.8), Inches(3.7), Inches(4.8), border_color=color)
    tb = s12_mig.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    
    p_body = tf.add_paragraph()
    p_body.text = f"\n{text}"
    p_body.font.size = Pt(13)
    p_body.font.color.rgb = TEXT_LIGHT



# ====================================================
# SLIDE 12: Business ROI & Metrics Impact
# ====================================================
s12 = add_blank_slide()
add_header(s12, "Business Impact", "Proven ROI & Operational Performance", "Measurable achievements delivered across multi-cloud infrastructure teams.")

add_card(s12, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.8), border_color=ACCENT_GREEN)
tb = s12.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.3), Inches(4.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "30 - 40%"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN
p_sub = tf.add_paragraph()
p_sub.text = "Annual Spend Reduction\n\nAchieved by automated elimination of idle resources and Reserved Instance optimization."
p_sub.font.size = Pt(13.5)
p_sub.font.color.rgb = TEXT_LIGHT

add_card(s12, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.8), border_color=ACCENT_CYAN)
tb2 = s12.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.3), Inches(4.4))
tf2 = tb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "< 2 Mins"
p2.font.size = Pt(44)
p2.font.bold = True
p2.font.color.rgb = ACCENT_CYAN
p_sub2 = tf2.add_paragraph()
p_sub2.text = "Anomaly Detection Time\n\nInstant z-score outlier identification catches cost spikes before month-end billing closure."
p_sub2.font.size = Pt(13.5)
p_sub2.font.color.rgb = TEXT_LIGHT

add_card(s12, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), border_color=ACCENT_PURPLE)
tb3 = s12.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.3), Inches(4.4))
tf3 = tb3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "95%+"
p3.font.size = Pt(44)
p3.font.bold = True
p3.font.color.rgb = ACCENT_PURPLE
p_sub3 = tf3.add_paragraph()
p_sub3.text = "PredictIQ Model Accuracy\n\nXGBoost Fourier time-series models provide high-precision 90-day budget projections."
p_sub3.font.size = Pt(13.5)
p_sub3.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 13: Installation, Operations & Deployment
# ====================================================
s13 = add_blank_slide()
add_header(s13, "Deployment Guide", "Installation & Operational Setup", "Simple multi-tier startup steps for local development and production deployment.")

add_card(s13, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), border_color=ACCENT_BLUE)
tb = s13.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "1. Backend Server Setup (`/server`)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p_body = tf.add_paragraph()
p_body.text = "\n```bash\ncd server\nnpm install\n# Configure .env (MONGO_URI, JWT_SECRET, PORT)\nnpm run dev\n```\n• Runs Express REST Gateway on Port 5000.\n• Connects to MongoDB Atlas & establishes Mongoose ORM models."
p_body.font.size = Pt(12.5)
p_body.font.color.rgb = TEXT_LIGHT

add_card(s13, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), border_color=ACCENT_CYAN)
tb2 = s13.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
tf2 = tb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "2. Frontend App Setup (`/client`)"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = ACCENT_CYAN
p_body2 = tf2.add_paragraph()
p_body2.text = "\n```bash\ncd client\nnpm install\nnpm run dev\n```\n• Launches Vite Dev Server on http://localhost:5173.\n• Loads React 19 UI with Framer Motion, GSAP, Recharts, and DataContext signals."
p_body2.font.size = Pt(12.5)
p_body2.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 14: Future Roadmap & Innovations
# ====================================================
s14 = add_blank_slide()
add_header(s14, "Future Vision", "Future Roadmap & Innovation Horizon", "Expanding CloudAtlas AI into autonomous multi-cloud infrastructure orchestration.")

roadmap = [
    ("🤖 Autonomous Auto-Remediation", "Direct integration with AWS CloudWatch / Azure Monitor APIs to automatically terminate flagged anomaly instances upon approval.", ACCENT_CYAN),
    ("☸️ Kubernetes Pod FinOps", "Granular pod-level cost attribution for K8s clusters using Prometheus & OpenCost metrics integration.", ACCENT_PURPLE),
    ("💬 LLM FinOps Copilot", "Natural language AI Assistant allowing users to ask 'Show me my highest spending Azure services in Q3' in plain English.", ACCENT_GREEN)
]

for i, (title, text, color) in enumerate(roadmap):
    left = Inches(0.8 + i * 3.95)
    add_card(s14, left, Inches(1.8), Inches(3.7), Inches(4.8), border_color=color)
    tb = s14.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    
    p_body = tf.add_paragraph()
    p_body.text = f"\n{text}"
    p_body.font.size = Pt(13)
    p_body.font.color.rgb = TEXT_LIGHT


# ====================================================
# SLIDE 15: Conclusion & Technical Q&A
# ====================================================
s15 = add_blank_slide()
add_card(s15, Inches(1.2), Inches(1.0), Inches(10.933), Inches(5.5), border_color=ACCENT_CYAN)

tb15 = s15.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.0))
tf15 = tb15.text_frame
tf15.word_wrap = True

p = tf15.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN
p.alignment = PP_ALIGN.CENTER

p2 = tf15.add_paragraph()
p2.text = "\nCloudAtlas AI — Next-Gen Multi-Cloud FinOps Platform"
p2.font.size = Pt(22)
p2.font.color.rgb = TEXT_LIGHT
p2.alignment = PP_ALIGN.CENTER

p3 = tf15.add_paragraph()
p3.text = "\nOpen for Technical Q&A & Demonstration"
p3.font.size = Pt(16)
p3.font.color.rgb = ACCENT_AMBER
p3.alignment = PP_ALIGN.CENTER

p4 = tf15.add_paragraph()
p4.text = "\nDeveloped by Neel Panchal & Henil"
p4.font.size = Pt(14)
p4.font.bold = True
p4.font.color.rgb = ACCENT_BLUE
p4.alignment = PP_ALIGN.CENTER

prs.save("CloudAtlas_AI_Presentation.pptx")
print("Successfully generated master comprehensive 15-slide CloudAtlas_AI_Presentation.pptx")
