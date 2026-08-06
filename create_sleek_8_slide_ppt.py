import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Modern Minimalist Dark Theme Palette (Stripe/Linear style)
BG_DARK = RGBColor(11, 15, 25)          # Deep Rich Dark Obsidian
TEXT_LIGHT = RGBColor(248, 250, 252)   # Pure Bright White
TEXT_MUTED = RGBColor(148, 163, 184)   # Soft Slate Gray
TEXT_DIM = RGBColor(100, 116, 139)     # Dimmed Slate
ACCENT_BLUE = RGBColor(56, 189, 248)   # Sky Blue
ACCENT_CYAN = RGBColor(45, 212, 191)   # Bright Cyan
ACCENT_PURPLE = RGBColor(168, 85, 247) # Neon Purple
ACCENT_GREEN = RGBColor(52, 211, 153)  # Mint Emerald
ACCENT_AMBER = RGBColor(251, 191, 36)  # Warm Gold
ACCENT_ROSE = RGBColor(244, 63, 94)    # Rose Red
LINE_COLOR = RGBColor(30, 41, 59)      # Subtle Divider Lines

def add_slide_bg(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    # Solid minimal obsidian background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    
    # Top thin accent line
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.02))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = LINE_COLOR
    top_line.line.fill.background()

    # Footer banner
    footer_tb = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.35))
    tf_f = footer_tb.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "CloudAtlas AI  •  Multi-Cloud FinOps & Predictive Governance Platform"
    p_f.font.size = Pt(10)
    p_f.font.color.rgb = TEXT_DIM
    p_f.font.name = 'Segoe UI'
    
    return slide

def add_header(slide, category, title):
    # Category / Kicker
    tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.733), Inches(0.35))
    tf_cat = tb_cat.text_frame
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_CYAN
    p_cat.font.name = 'Segoe UI'
    
    # Title
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.65))
    tf_title = tb_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT
    p_title.font.name = 'Segoe UI'

# ====================================================
# SLIDE 1: Title Slide — Minimalist Modern Hero (No Boxes)
# ====================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = BG_DARK
bg1.line.fill.background()

# Glowing Top Decorative Line
tline = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05))
tline.fill.solid()
tline.fill.fore_color.rgb = ACCENT_CYAN
tline.line.fill.background()

tb1 = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "CloudAtlas AI"
p1.font.size = Pt(54)
p1.font.bold = True
p1.font.color.rgb = TEXT_LIGHT
p1.font.name = 'Segoe UI'

p2 = tf1.add_paragraph()
p2.text = "Next-Generation Multi-Cloud FinOps & Predictive Governance"
p2.font.size = Pt(24)
p2.font.color.rgb = ACCENT_CYAN
p2.font.name = 'Segoe UI'

p3 = tf1.add_paragraph()
p3.text = "\nUnifying AWS, Azure & GCP cost analytics with 90-day XGBoost time-series forecasting\nand unsupervised One-Class SVM anomaly detection."
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED
p3.font.name = 'Segoe UI'

# Clean Stat Row (Minimal, No Boxes, Just Text Counters)
stats = [
    ("90 Days", "XGBoost Time-Series Forecast", ACCENT_CYAN),
    ("35%+", "Annual Cloud Cost Reduction", ACCENT_GREEN),
    ("Zero-Trust", "JWT & Gatekeeper Shield", ACCENT_PURPLE),
    ("Multi-Cloud", "AWS • Azure • GCP", ACCENT_AMBER)
]

for i, (val, label, col) in enumerate(stats):
    sx = Inches(0.8) + i * Inches(2.95)
    tb_s = s1.shapes.add_textbox(sx, Inches(5.2), Inches(2.7), Inches(1.4))
    tf_s = tb_s.text_frame
    tf_s.word_wrap = True
    
    ps1 = tf_s.paragraphs[0]
    ps1.text = val
    ps1.font.size = Pt(28)
    ps1.font.bold = True
    ps1.font.color.rgb = col
    ps1.font.name = 'Segoe UI'
    
    ps2 = tf_s.add_paragraph()
    ps2.text = label
    ps2.font.size = Pt(11)
    ps2.font.color.rgb = TEXT_MUTED
    ps2.font.name = 'Segoe UI'

# ====================================================
# SLIDE 2: Problem Statement & Strategic Vision (Split Layout, No Boxes)
# ====================================================
s2 = add_slide_bg(prs)
add_header(s2, "01. Executive Problem & Vision", "The Cloud Cost Crisis & The CloudAtlas Advantage")

# Left Column: Big Highlight Statement
tb_left = s2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.8))
tf_left = tb_left.text_frame
tf_left.word_wrap = True

pl1 = tf_left.paragraphs[0]
pl1.text = "Modern cloud spend is opaque, volatile, and reactive."
pl1.font.size = Pt(28)
pl1.font.bold = True
pl1.font.color.rgb = ACCENT_ROSE
pl1.font.name = 'Segoe UI'

pl2 = tf_left.add_paragraph()
pl2.text = "\nEnterprises waste billions annually on orphaned storage, unmonitored serverless loops, and inaccurate spend forecasts."
pl2.font.size = Pt(14)
pl2.font.color.rgb = TEXT_MUTED
pl2.font.name = 'Segoe UI'

# Middle Divider Line
div1 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.3), Inches(1.8), Inches(0.02), Inches(4.8))
div1.fill.solid()
div1.fill.fore_color.rgb = LINE_COLOR
div1.line.fill.background()

# Right Column: 3 Clean Bullet Highlights with Colored Accent Indicators
rights = [
    ("⚡ Raw Billing Opacity", "Cloud providers export massive CSV logs with millions of unaggregated rows. FinOps teams lack unified cross-cloud visibility.", ACCENT_AMBER),
    ("💥 Unmitigated Cost Spikes", "Orphaned EBS disks and runaway compute jobs trigger budget overruns before traditional monthly invoices arrive.", ACCENT_ROSE),
    ("🔮 Proactive Governance Solution", "CloudAtlas AI unifies AWS, Azure, and GCP billing logs into an automated machine learning engine for instant forecasting and outlier detection.", ACCENT_CYAN)
]

for i, (title, desc, col) in enumerate(rights):
    ry = Inches(1.8) + i * Inches(1.6)
    
    # Bullet dot indicator
    dot = s2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.6), ry + Inches(0.08), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    
    tb_r = s2.shapes.add_textbox(Inches(6.85), ry, Inches(5.6), Inches(1.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    pr1 = tf_r.paragraphs[0]
    pr1.text = title
    pr1.font.size = Pt(16)
    pr1.font.bold = True
    pr1.font.color.rgb = TEXT_LIGHT
    pr1.font.name = 'Segoe UI'
    
    pr2 = tf_r.add_paragraph()
    pr2.text = desc
    pr2.font.size = Pt(12)
    pr2.font.color.rgb = TEXT_MUTED
    pr2.font.name = 'Segoe UI'

# ====================================================
# SLIDE 3: Technical Microservices Architecture (Flow Pipeline Layout)
# ====================================================
s3 = add_slide_bg(prs)
add_header(s3, "02. Platform Architecture", "Clean 3-Tier Enterprise Microservice Stack")

architectures = [
    ("1. Frontend Client", "React 19 • Vite 6 • Framer Motion", [
        "Spring-animated interactive UI components",
        "GSAP ScrollTrigger horizontal tech showcase",
        "Recharts interactive cost distribution charts",
        "Tailwind & custom dark glassmorphism design"
    ], ACCENT_CYAN),
    ("2. Backend Gateway", "Node.js • Express REST • MongoDB", [
        "Mongoose ORM schema architecture",
        "Zero-Trust JWT authentication guards",
        "Multer streaming multi-part CSV ingest",
        "Express rate limiters & SOC2 audit logs"
    ], ACCENT_PURPLE),
    ("3. Machine Learning", "Python • Fast Pipeline • Scikit-Learn", [
        "PredictIQ 90-day XGBoost spend regressor",
        "Unsupervised One-Class SVM anomaly shield",
        "Fourier Sine seasonality mathematical modeling",
        "What-If workload scale simulation engine"
    ], ACCENT_GREEN)
]

for i, (tier_title, tech_stack, bullets, col) in enumerate(architectures):
    cx = Inches(0.8) + i * Inches(3.95)
    
    # Left vertical colored bar instead of full box
    bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, Inches(1.8), Inches(0.06), Inches(4.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    
    tb = s3.shapes.add_textbox(cx + Inches(0.2), Inches(1.75), Inches(3.5), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = tier_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = col
    p.font.name = 'Segoe UI'
    
    pt = tf.add_paragraph()
    pt.text = tech_stack
    pt.font.size = Pt(12)
    pt.font.bold = True
    pt.font.color.rgb = TEXT_LIGHT
    pt.font.name = 'Segoe UI'
    
    tf.add_paragraph().text = ""
    
    for bullet in bullets:
        pb = tf.add_paragraph()
        pb.text = f"• {bullet}"
        pb.font.size = Pt(11)
        pb.font.color.rgb = TEXT_MUTED
        pb.font.name = 'Segoe UI'

# ====================================================
# SLIDE 4: Core Feature Showcase (Clean Minimal Layout)
# ====================================================
s4 = add_slide_bg(prs)
add_header(s4, "03. Core Feature Showcase", "Interactive Modules Designed for Seamless FinOps Experience")

modules = [
    ("🎨 FinOps Landing Sandbox", "Interactive 5-tab live showcase featuring ROI spend sliders calculating immediate annual cost savings ($5k - $500k+).", ACCENT_CYAN),
    ("📊 Executive Dashboard", "Unified multi-cloud views aggregating monthly spend, provider distribution (AWS, Azure, GCP), and top expensive services.", ACCENT_BLUE),
    ("🔐 6-Digit OTP Auth & Security Gate", "Zero-Trust authentication flow with 2-minute OTP countdown and Master Gatekeeper passcode challenge (`HenilNeelProject`).", ACCENT_PURPLE),
    ("📂 Automated CSV Ingestion Engine", "High-performance parser standardizing multi-cloud provider billing logs with instant global data context state sync.", ACCENT_GREEN)
]

for i, (m_title, m_desc, col) in enumerate(modules):
    row = i // 2
    c = i % 2
    mx = Inches(0.8) + c * Inches(5.95)
    my = Inches(1.8) + row * Inches(2.5)
    
    # Left accent vertical border line
    bar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, mx, my, Inches(0.06), Inches(2.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    
    tb = s4.shapes.add_textbox(mx + Inches(0.2), my, Inches(5.4), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = m_title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = col
    p.font.name = 'Segoe UI'
    
    p2 = tf.add_paragraph()
    p2.text = m_desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_LIGHT
    p2.font.name = 'Segoe UI'

# ====================================================
# SLIDE 5: Machine Learning Engine (PredictIQ & SVM)
# ====================================================
s5 = add_slide_bg(prs)
add_header(s5, "04. Machine Learning Engine", "PredictIQ 90-Day Forecast & One-Class SVM Anomaly Shield")

# Left Column: PredictIQ
tb_m1 = s5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
tf_m1 = tb_m1.text_frame
tf_m1.word_wrap = True

p1 = tf_m1.paragraphs[0]
p1.text = "🔮 PredictIQ 90-Day Forecast Engine"
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = ACCENT_CYAN
p1.font.name = 'Segoe UI'

p_sub1 = tf_m1.add_paragraph()
p_sub1.text = "XGBoost Time-Series Regressor with Fourier Seasonality\n"
p_sub1.font.size = Pt(12)
p_sub1.font.bold = True
p_sub1.font.color.rgb = TEXT_LIGHT

pts1 = [
    "Predicts daily and monthly cost trajectories up to 90 days in advance.",
    "Combines linear trend slope with cyclical Fourier sine seasonality:",
    "  Seasonality(t) = sin(t * 0.4) * 0.09 * mean_cost",
    "Generates dynamic upper and lower confidence envelopes.",
    "Eliminates budget end-of-quarter surprises for CFOs."
]
for pt in pts1:
    p_pt = tf_m1.add_paragraph()
    p_pt.text = f"• {pt}" if not pt.startswith("  ") else pt
    p_pt.font.size = Pt(12)
    p_pt.font.color.rgb = TEXT_MUTED if not pt.startswith("  ") else ACCENT_AMBER

# Divider Line
div2 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.8), Inches(0.02), Inches(4.8))
div2.fill.solid()
div2.fill.fore_color.rgb = LINE_COLOR
div2.line.fill.background()

# Right Column: One-Class SVM
tb_m2 = s5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
tf_m2 = tb_m2.text_frame
tf_m2.word_wrap = True

p2 = tf_m2.paragraphs[0]
p2.text = "🛡️ One-Class SVM Anomaly Shield"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = ACCENT_ROSE
p2.font.name = 'Segoe UI'

p_sub2 = tf_m2.add_paragraph()
p_sub2.text = "Unsupervised Outlier Detection on Billing Logs\n"
p_sub2.font.size = Pt(12)
p_sub2.font.bold = True
p_sub2.font.color.rgb = TEXT_LIGHT

pts2 = [
    "Operates seamlessly on unlabelled multi-cloud expenditure datasets.",
    "Fits an n-dimensional hypersphere decision boundary around normal metrics.",
    "Flags statistical outliers where cost variance threshold z >= 1.6.",
    "Automatically suggests remediation (e.g., terminate orphan EBS volumes).",
    "Stops runaway cloud spend before invoices are generated."
]
for pt in pts2:
    p_pt = tf_m2.add_paragraph()
    p_pt.text = f"• {pt}"
    p_pt.font.size = Pt(12)
    p_pt.font.color.rgb = TEXT_MUTED

# ====================================================
# SLIDE 6: Interactive FinOps Sandbox & What-If Simulator
# ====================================================
s6 = add_slide_bg(prs)
add_header(s6, "05. FinOps Intelligence", "What-If Workload Simulator & Multi-Cloud Optimization")

sandbox_items = [
    ("🎛️ Interactive Resource Sliders", "Real-time sliders for CPU %, Memory %, Storage GB, and Network Egress GB dynamically recalculating 90-day forecast curves.", ACCENT_CYAN),
    ("💳 Reserved Instance (RI) Advisor", "Instant financial comparison evaluating On-Demand pay-as-you-go costs versus 1-Year / 3-Year Reserved Instance commitments.", ACCENT_GREEN),
    ("☁️ Multi-Cloud Provider Migration", "Granular workload placement calculator comparing AWS, Azure, and GCP region pricing to identify maximum cost efficiency.", ACCENT_AMBER),
    ("💡 AI Actionable Recommendations", "Context-aware optimization tips highlighting unattached volumes, idle load balancers, and over-provisioned instance types.", ACCENT_PURPLE)
]

for i, (stitle, sdesc, col) in enumerate(sandbox_items):
    row = i // 2
    c = i % 2
    sx = Inches(0.8) + c * Inches(5.95)
    sy = Inches(1.8) + row * Inches(2.5)
    
    # Left accent vertical bar
    bar = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, sy, Inches(0.06), Inches(2.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    
    tb = s6.shapes.add_textbox(sx + Inches(0.2), sy, Inches(5.4), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = stitle
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = col
    p.font.name = 'Segoe UI'
    
    p2 = tf.add_paragraph()
    p2.text = sdesc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_LIGHT
    p2.font.name = 'Segoe UI'

# ====================================================
# SLIDE 7: Security, Compliance & Data Governance
# ====================================================
s7 = add_slide_bg(prs)
add_header(s7, "06. Security & Compliance", "Zero-Trust Security & Enterprise Governance Shield")

sec_features = [
    ("🔒 Zero-Trust JWT Authentication", "All API routes protected with cryptographically verified JSON Web Tokens and granular user authorization scopes.", ACCENT_PURPLE),
    ("🔑 Super Admin Passcode Challenge", "Restricted administrative portal (`/admin/login`) secured behind a dual-layer Master Gatekeeper Passcode (`HenilNeelProject`).", ACCENT_CYAN),
    ("📋 SOC2 Audit Logging Pipeline", "Complete audit trail tracking every authentication attempt, billing CSV upload, and dataset deletion event.", ACCENT_AMBER),
    ("🛡️ Rate Limiting & DDoS Defense", "Express rate-limiters preventing brute-force login attacks, credential stuffing, and unauthenticated API scraping.", ACCENT_GREEN)
]

for i, (sec_title, sec_desc, col) in enumerate(sec_features):
    row = i // 2
    c = i % 2
    sec_x = Inches(0.8) + c * Inches(5.95)
    sec_y = Inches(1.8) + row * Inches(2.5)
    
    # Horizontal top colored bar
    hbar = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, sec_x, sec_y, Inches(5.5), Inches(0.04))
    hbar.fill.solid()
    hbar.fill.fore_color.rgb = col
    hbar.line.fill.background()
    
    tb = s7.shapes.add_textbox(sec_x, sec_y + Inches(0.1), Inches(5.5), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = sec_title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = 'Segoe UI'
    
    p2 = tf.add_paragraph()
    p2.text = sec_desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = 'Segoe UI'

# ====================================================
# SLIDE 8: Enterprise Value & Future Roadmap
# ====================================================
s8 = add_slide_bg(prs)
add_header(s8, "07. ROI & Strategic Roadmap", "Measurable Enterprise Impact & Product Expansion")

# Left Column: Big ROI Callouts (Clean Numbers)
tb_r1 = s8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
tf_r1 = tb_r1.text_frame
tf_r1.word_wrap = True

pr1 = tf_r1.paragraphs[0]
pr1.text = "Business Impact"
pr1.font.size = Pt(20)
pr1.font.bold = True
pr1.font.color.rgb = ACCENT_GREEN
pr1.font.name = 'Segoe UI'

roi_stats = [
    ("35%+", "Annual Cloud Cost Savings via automated idle resource discovery"),
    ("Zero", "Budget Spikes with real-time One-Class SVM anomaly alerts"),
    ("90 Days", "Proactive financial forecasting visibility for C-suite teams")
]

for stat_val, stat_lbl in roi_stats:
    p_v = tf_r1.add_paragraph()
    p_v.text = stat_val
    p_v.font.size = Pt(26)
    p_v.font.bold = True
    p_v.font.color.rgb = ACCENT_CYAN
    
    p_l = tf_r1.add_paragraph()
    p_l.text = stat_lbl + "\n"
    p_l.font.size = Pt(11)
    p_l.font.color.rgb = TEXT_MUTED

# Divider Line
div3 = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.8), Inches(0.02), Inches(4.8))
div3.fill.solid()
div3.fill.fore_color.rgb = LINE_COLOR
div3.line.fill.background()

# Right Column: Clean Timeline Roadmap
tb_r2 = s8.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
tf_r2 = tb_r2.text_frame
tf_r2.word_wrap = True

pr2 = tf_r2.paragraphs[0]
pr2.text = "Strategic Future Roadmap"
pr2.font.size = Pt(20)
pr2.font.bold = True
pr2.font.color.rgb = ACCENT_PURPLE
pr2.font.name = 'Segoe UI'

phases = [
    ("Phase 1 • Auto-Remediation", "Automated Cloud Webhooks terminating orphaned resources without manual intervention."),
    ("Phase 2 • Kubernetes (k8s) Granularity", "Pod and namespace level cost allocation across EKS, AKS, and GKE clusters."),
    ("Phase 3 • Enterprise Multi-Tenant RBAC", "Role-based access controls with dedicated department budget quotas."),
    ("Phase 4 • LLM FinOps Copilot", "Natural language querying engine ('Why did AWS compute costs jump yesterday?').")
]

for title, desc in phases:
    p_t = tf_r2.add_paragraph()
    p_t.text = title
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_LIGHT
    
    p_d = tf_r2.add_paragraph()
    p_d.text = desc + "\n"
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = TEXT_MUTED

# Save
output_path = r"c:\Users\NEEL\Downloads\OneDrive\Desktop\Today\CloudAtlas_AI_Sleek_8_Slide_Presentation.pptx"
prs.save(output_path)
print(f"SUCCESS: Sleek Minimalist PPT created at {output_path}")
